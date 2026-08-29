#!/usr/bin/env python3
"""
PPTX Deck Renderer — builds a PowerPoint deck directly from PresentationModel.

Row f507034e. Marp's ``--pptx`` export rasterizes every slide to a PNG, so
finished decks carried zero selectable text. Rick ruled python-pptx (2026-08-16):
build the deck from the structured model instead, so every slide gets REAL,
selectable, screen-reader-visible text runs — not text painted over an image.

This renderer is a pure, synchronous transformation: PresentationModel + theme
config dict + a visuals directory in, a ``.pptx`` file on disk out. No LLM, no
async, no network. It mirrors MarpTextRenderer's slide-type dispatch so the two
paths stay recognizably the same shape.

Genuine content images (matplotlib charts, generated images, video frames) that
Phase 7 rendered into ``visuals/`` as raster PNGs are embedded as pictures.
SVG diagrams (d2) and Marp-native mermaid blocks are a known v1 gap — those
slides keep their text; the diagram is deferred to a follow-up (needs a
rasterizer / mmdc). See src/rnd/v0.2.0/2026.08.16-pptx-text-layer-python-pptx.md.
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# =============================================================================
# Layout constants (16:9, matching Marp's default aspect ratio)
# =============================================================================

_SLIDE_W_IN = 13.333
_SLIDE_H_IN = 7.5
_MARGIN_IN  = 0.6

# Raster visual filenames Phase 7 writes, keyed by placeholder index NNN.
# SVG (d2) is intentionally absent — python-pptx cannot embed it (v1 gap).
_VISUAL_FILENAME_TEMPLATES = [
    "chart-{idx:03d}.png",        # MatplotlibRenderer
    "image-{idx:03d}.png",        # NanoBananaRenderer
    "video-{idx:03d}-frame.png",  # VeoRenderer
]

# Default theme colors, used when theme_config omits a key. Mirrors
# MarpTextRenderer.DEFAULT_THEME_CONFIG so the two paths agree on defaults.
_DEFAULT_COLORS = {
    "primary"    : "2563EB",
    "secondary"  : "1E40AF",
    "accent"     : "F59E0B",
    "background" : "FFFFFF",
    "text"       : "1F2937",
}

# Slide types that render as a centered "lead" layout (title + subtitle only).
_LEAD_TYPES = { "title", "section_divider" }


class PptxDeckRenderer:
    """
    Stateless builder that converts a PresentationModel to a .pptx on disk.

    All methods are @staticmethod — no constructor, no instance state. The one
    public entry point is build().

    Requires:
        - presentation is a PresentationModel instance
        - theme_config is a dict matching the theme YAML schema
        - visuals_dir is a directory path (may not exist — treated as no visuals)
        - output_path is a writable .pptx file path

    Ensures:
        - Writes a .pptx to output_path whose slides carry real text runs
        - Returns output_path
    """

    # =========================================================================
    # Public entry point
    # =========================================================================

    @staticmethod
    def build( presentation, theme_config, visuals_dir, output_path ):
        """
        Build a .pptx from a PresentationModel and write it to output_path.

        Requires:
            - presentation has a .slides list of SlideModel
            - theme_config is a dict with an optional "theme" -> "colors" map
            - visuals_dir is a str path (existence not required)
            - output_path is a writable file path; parent dir is created

        Ensures:
            - Every slide gets native text (title/subtitle/bullets) + real notes
            - Non-text_only slides embed their raster visual if one exists on disk
            - Returns output_path

        Returns:
            str: output_path
        """
        colors = PptxDeckRenderer._resolve_colors( theme_config )

        prs = Presentation()
        prs.slide_width  = Inches( _SLIDE_W_IN )
        prs.slide_height = Inches( _SLIDE_H_IN )
        blank_layout     = prs.slide_layouts[ 6 ]

        visual_idx = 0
        for slide_model in presentation.slides:
            slide = prs.slides.add_slide( blank_layout )
            PptxDeckRenderer._paint_background( slide, colors[ "background" ] )

            # Resolve the visual file for this slide (only non-text_only slides
            # consume a placeholder index — matching Phase 7's enumeration).
            visual_path = None
            if slide_model.visual_type != "text_only":
                visual_path = PptxDeckRenderer._find_visual_file( visuals_dir, visual_idx )
                visual_idx += 1

            if slide_model.type in _LEAD_TYPES:
                PptxDeckRenderer._render_lead_slide( slide, slide_model, presentation, colors )
            else:
                PptxDeckRenderer._render_content_slide( slide, slide_model, colors, visual_path )

            PptxDeckRenderer._add_notes( slide, slide_model.presenter_notes )

        os.makedirs( os.path.dirname( os.path.abspath( output_path ) ), exist_ok=True )
        prs.save( output_path )
        return output_path

    # =========================================================================
    # Theme
    # =========================================================================

    @staticmethod
    def _resolve_colors( theme_config ):
        """
        Resolve theme colors to RGBColor objects, falling back to defaults.

        Requires:
            - theme_config is a dict (may be empty or missing keys)

        Ensures:
            - returns a dict mapping primary/secondary/accent/background/text to
              RGBColor, using _DEFAULT_COLORS for any missing entry
        """
        theme  = theme_config.get( "theme", {} ) if isinstance( theme_config, dict ) else {}
        colors = theme.get( "colors", {} ) if isinstance( theme, dict ) else {}

        resolved = {}
        for key, default_hex in _DEFAULT_COLORS.items():
            resolved[ key ] = PptxDeckRenderer._hex_to_rgb( colors.get( key, default_hex ) )
        return resolved

    @staticmethod
    def _hex_to_rgb( hex_str ):
        """
        Convert a hex color string to an RGBColor, tolerating a leading '#'.

        Requires:
            - hex_str is a 6-hex-digit string, optionally prefixed with '#'

        Ensures:
            - returns an RGBColor for a valid 6-hex-digit value
            - falls back to the text default for anything unparseable
        """
        cleaned = hex_str.lstrip( "#" ) if isinstance( hex_str, str ) else ""
        try:
            return RGBColor.from_string( cleaned )
        except ValueError:
            return RGBColor.from_string( _DEFAULT_COLORS[ "text" ] )

    @staticmethod
    def _paint_background( slide, bg_rgb ):
        """Set the slide's background to a solid theme color."""
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_rgb

    # =========================================================================
    # Visual file resolution
    # =========================================================================

    @staticmethod
    def _find_visual_file( visuals_dir, visual_idx ):
        """
        Find the raster visual file for a placeholder index, if one exists.

        Requires:
            - visuals_dir is a str path (may not exist)
            - visual_idx is a non-negative int

        Ensures:
            - returns the absolute path of the first matching raster file that
              exists on disk (chart / image / video-frame PNG), or None
        """
        for template in _VISUAL_FILENAME_TEMPLATES:
            candidate = os.path.join( visuals_dir, template.format( idx=visual_idx ) )
            if os.path.isfile( candidate ):
                return candidate
        return None

    # =========================================================================
    # Slide renderers
    # =========================================================================

    @staticmethod
    def _render_lead_slide( slide, slide_model, presentation, colors ):
        """
        Render a centered lead slide (title / section_divider): title + subtitle
        + optional speaker | date line.
        """
        box = slide.shapes.add_textbox(
            Inches( _MARGIN_IN ), Inches( 2.2 ),
            Inches( _SLIDE_W_IN - 2 * _MARGIN_IN ), Inches( 3.1 ),
        )
        tf = box.text_frame
        tf.word_wrap = True

        PptxDeckRenderer._set_paragraph(
            tf.paragraphs[ 0 ], slide_model.title, Pt( 40 ), colors[ "primary" ],
            bold=True, align=PP_ALIGN.CENTER,
        )

        if slide_model.subtitle:
            PptxDeckRenderer._set_paragraph(
                tf.add_paragraph(), slide_model.subtitle, Pt( 24 ), colors[ "secondary" ],
                align=PP_ALIGN.CENTER,
            )

        speaker_date = " | ".join(
            part for part in ( presentation.speaker, presentation.date ) if part
        )
        if speaker_date:
            PptxDeckRenderer._set_paragraph(
                tf.add_paragraph(), speaker_date, Pt( 18 ), colors[ "text" ],
                align=PP_ALIGN.CENTER,
            )

    @staticmethod
    def _render_content_slide( slide, slide_model, colors, visual_path ):
        """
        Render a content slide: title, optional subtitle, bullets, and — when a
        raster visual exists — the picture beside the text.
        """
        # Title band across the top.
        title_box = slide.shapes.add_textbox(
            Inches( _MARGIN_IN ), Inches( 0.4 ),
            Inches( _SLIDE_W_IN - 2 * _MARGIN_IN ), Inches( 1.2 ),
        )
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        PptxDeckRenderer._set_paragraph(
            title_tf.paragraphs[ 0 ], slide_model.title, Pt( 32 ), colors[ "primary" ], bold=True,
        )
        if slide_model.subtitle:
            PptxDeckRenderer._set_paragraph(
                title_tf.add_paragraph(), slide_model.subtitle, Pt( 20 ), colors[ "secondary" ],
            )

        # Body region below the title. If a visual exists, text takes the left
        # half and the picture the right half; otherwise text spans full width.
        body_top    = 1.8
        body_height = _SLIDE_H_IN - body_top - _MARGIN_IN
        full_width  = _SLIDE_W_IN - 2 * _MARGIN_IN
        text_width  = full_width / 2 if visual_path else full_width

        if slide_model.content_bullets:
            bullets_box = slide.shapes.add_textbox(
                Inches( _MARGIN_IN ), Inches( body_top ),
                Inches( text_width ), Inches( body_height ),
            )
            bullets_tf = bullets_box.text_frame
            bullets_tf.word_wrap = True
            for i, bullet in enumerate( slide_model.content_bullets ):
                para = bullets_tf.paragraphs[ 0 ] if i == 0 else bullets_tf.add_paragraph()
                PptxDeckRenderer._set_paragraph(
                    para, f"• {bullet}", Pt( 20 ), colors[ "text" ],
                )

        if visual_path:
            PptxDeckRenderer._add_picture(
                slide, visual_path,
                left_in=_MARGIN_IN + text_width, top_in=body_top,
                width_in=text_width, height_in=body_height,
            )

    @staticmethod
    def _add_picture( slide, image_path, left_in, top_in, width_in, height_in ):
        """
        Add a picture constrained to a box, scaled to fit while preserving aspect.
        """
        pic = slide.shapes.add_picture( image_path, Inches( left_in ), Inches( top_in ) )

        box_w = Inches( width_in )
        box_h = Inches( height_in )
        scale = min( box_w / pic.width, box_h / pic.height, 1 )

        pic.width  = int( pic.width * scale )
        pic.height = int( pic.height * scale )
        # Center within the box.
        pic.left = Inches( left_in ) + ( box_w - pic.width ) // 2
        pic.top  = Inches( top_in ) + ( box_h - pic.height ) // 2

    # =========================================================================
    # Presenter notes
    # =========================================================================

    @staticmethod
    def _add_notes( slide, notes ):
        """
        Populate the slide's real notes slide from PresenterNotes.

        Requires:
            - slide is a python-pptx slide
            - notes is a PresenterNotes instance or None

        Ensures:
            - writes a notes slide only when there is note content; a slide with
              no notes gets no notes slide
        """
        text = PptxDeckRenderer._format_notes( notes )
        if text:
            slide.notes_slide.notes_text_frame.text = text

    @staticmethod
    def _format_notes( notes ):
        """
        Format PresenterNotes into a plain-text notes string (empty if none).
        """
        if notes is None:
            return ""

        lines = []
        if notes.transition:
            lines.append( f"Transition: {notes.transition}" )
        if notes.talking_points:
            lines.append( "Talking points:" )
            for point in notes.talking_points:
                lines.append( f"- {point}" )
        if notes.timing_seconds and notes.timing_seconds > 0:
            lines.append( f"Timing: {notes.timing_seconds}s" )
        if notes.emphasis:
            lines.append( f"Emphasis: {notes.emphasis}" )

        return "\n".join( lines )

    # =========================================================================
    # Text helper
    # =========================================================================

    @staticmethod
    def _set_paragraph( paragraph, text, size, rgb, bold=False, align=PP_ALIGN.LEFT ):
        """
        Write one run of styled text into a paragraph.

        Requires:
            - paragraph is a python-pptx paragraph
            - text is a string, size a Pt length, rgb an RGBColor

        Ensures:
            - the paragraph carries a single styled run with the given text
        """
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = text
        run.font.size = size
        run.font.bold = bold
        run.font.color.rgb = rgb


# =============================================================================
# Smoke Test
# =============================================================================

def quick_smoke_test():
    """Quick smoke test for PptxDeckRenderer."""
    import sys
    import tempfile

    lupin_root = os.environ.get( "LUPIN_ROOT" )
    if lupin_root:
        src_path = os.path.join( lupin_root, "src" )
        if src_path not in sys.path:
            sys.path.insert( 0, src_path )

    from cosa.agents.presentation_generator.state import (
        PresentationModel, SlideModel, PresenterNotes
    )
    from cosa.agents.presentation_generator.deck_verdict import verify_presentation_deck

    print( "=" * 60 )
    print( "PptxDeckRenderer Smoke Test" )
    print( "=" * 60 )

    try:
        presentation = PresentationModel(
            title   = "Owning Our Layout",
            speaker = "Tiffany",
            date    = "2026-08-16",
            slides  = [
                SlideModel(
                    number=1, arc_position="opening", type="title",
                    title="Owning Our Layout", subtitle="Real text, at last",
                    presenter_notes=PresenterNotes( talking_points=[ "Welcome" ], timing_seconds=30 ),
                ),
                SlideModel(
                    number=2, arc_position="body", type="content",
                    title="Why This Matters", content_bullets=[ "Selectable", "Accessible", "Editable" ],
                    presenter_notes=PresenterNotes( transition="So...", talking_points=[ "Explain" ] ),
                ),
            ],
        )

        out_dir = tempfile.mkdtemp( prefix="pptx_smoke_" )
        out_path = os.path.join( out_dir, "smoke.pptx" )

        print( "\nTest 1: Build deck..." )
        result_path = PptxDeckRenderer.build( presentation, {}, out_dir, out_path )
        assert os.path.isfile( result_path ), "deck not written"
        print( "  PASS" )

        print( "\nTest 2: Deck has a real text layer..." )
        verdict = verify_presentation_deck( result_path )
        assert verdict, f"verdict failed: {verdict.reason}"
        assert verdict.text_run_count > 0, "no text runs"
        assert verdict.slide_count == 2, f"slide count {verdict.slide_count}"
        print( f"  PASS ({verdict.text_run_count} text runs across {verdict.slide_count} slides)" )

        print( f"\n{'=' * 60}" )
        print( "All PptxDeckRenderer smoke tests passed" )
        print( f"{'=' * 60}" )

    except Exception as e:
        print( f"\nSmoke test FAILED: {e}" )
        import traceback
        traceback.print_exc()


if __name__ == "__main__":
    quick_smoke_test()
