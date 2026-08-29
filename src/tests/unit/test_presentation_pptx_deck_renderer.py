"""
Unit tests for PptxDeckRenderer (row f507034e).

The renderer builds a PowerPoint deck from a PresentationModel so slides carry
REAL, selectable text runs instead of Marp's rasterized images. These tests
assert the text layer exists, notes are populated, raster visuals embed, theme
colors resolve, and every slide-type / branch is exercised.

Venue: :7999 (pure, tmp fixtures, sub-second, no server, no LLM) — meets all
three §TESTING VENUES criteria.
"""

import os

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

from cosa.agents.presentation_generator.state import (
    PresentationModel, SlideModel, PresenterNotes,
)
from cosa.agents.presentation_generator.renderers.pptx_deck_renderer import (
    PptxDeckRenderer,
)
from cosa.agents.presentation_generator.deck_verdict import verify_presentation_deck


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _png( path, size=( 400, 300 ), color=( 30, 60, 200 ) ):
    Image.new( "RGB", size, color ).save( str( path ) )
    return str( path )


def _slide_text_runs( pptx_path ):
    """Total non-empty <a:t> runs across the built deck's slides (native read)."""
    prs = Presentation( pptx_path )
    total = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text:
                            total += 1
    return total


# ---------------------------------------------------------------------------
# _hex_to_rgb
# ---------------------------------------------------------------------------

def test_hex_to_rgb_plain():
    assert PptxDeckRenderer._hex_to_rgb( "2563EB" ) == RGBColor.from_string( "2563EB" )


def test_hex_to_rgb_with_hash():
    assert PptxDeckRenderer._hex_to_rgb( "#F59E0B" ) == RGBColor.from_string( "F59E0B" )


def test_hex_to_rgb_non_string_falls_back():
    # Non-str → cleaned "" → ValueError → text default.
    assert PptxDeckRenderer._hex_to_rgb( None ) == RGBColor.from_string( "1F2937" )


def test_hex_to_rgb_bad_value_falls_back():
    assert PptxDeckRenderer._hex_to_rgb( "nothex" ) == RGBColor.from_string( "1F2937" )


# ---------------------------------------------------------------------------
# _resolve_colors
# ---------------------------------------------------------------------------

def test_resolve_colors_all_defaults_when_empty():
    colors = PptxDeckRenderer._resolve_colors( {} )
    assert colors[ "primary" ]    == RGBColor.from_string( "2563EB" )
    assert colors[ "background" ] == RGBColor.from_string( "FFFFFF" )


def test_resolve_colors_non_dict_theme_config():
    colors = PptxDeckRenderer._resolve_colors( None )
    assert colors[ "text" ] == RGBColor.from_string( "1F2937" )


def test_resolve_colors_theme_not_a_dict():
    colors = PptxDeckRenderer._resolve_colors( { "theme": "not-a-dict" } )
    assert colors[ "primary" ] == RGBColor.from_string( "2563EB" )


def test_resolve_colors_honors_overrides():
    cfg = { "theme": { "colors": { "primary": "#123456" } } }
    colors = PptxDeckRenderer._resolve_colors( cfg )
    assert colors[ "primary" ]   == RGBColor.from_string( "123456" )
    assert colors[ "secondary" ] == RGBColor.from_string( "1E40AF" )   # default kept


# ---------------------------------------------------------------------------
# _find_visual_file
# ---------------------------------------------------------------------------

def test_find_visual_file_chart_match( tmp_path ):
    _png( tmp_path / "chart-000.png" )
    found = PptxDeckRenderer._find_visual_file( str( tmp_path ), 0 )
    assert found and found.endswith( "chart-000.png" )


def test_find_visual_file_video_frame_match( tmp_path ):
    # Only the third template matches — exercises iterating past non-matches.
    _png( tmp_path / "video-002-frame.png" )
    found = PptxDeckRenderer._find_visual_file( str( tmp_path ), 2 )
    assert found and found.endswith( "video-002-frame.png" )


def test_find_visual_file_none_when_missing( tmp_path ):
    assert PptxDeckRenderer._find_visual_file( str( tmp_path ), 7 ) is None


# ---------------------------------------------------------------------------
# _format_notes
# ---------------------------------------------------------------------------

def test_format_notes_none_is_empty():
    assert PptxDeckRenderer._format_notes( None ) == ""


def test_format_notes_all_fields():
    notes = PresenterNotes(
        transition="So then...", talking_points=[ "A", "B" ],
        timing_seconds=90, emphasis="Pause here",
    )
    text = PptxDeckRenderer._format_notes( notes )
    assert "Transition: So then..." in text
    assert "Talking points:" in text and "- A" in text and "- B" in text
    assert "Timing: 90s" in text
    assert "Emphasis: Pause here" in text


def test_format_notes_zero_timing_and_no_fields_is_empty():
    # timing 0 skips the timing line; nothing else set → empty string.
    assert PptxDeckRenderer._format_notes( PresenterNotes( timing_seconds=0 ) ) == ""


# ---------------------------------------------------------------------------
# _set_paragraph
# ---------------------------------------------------------------------------

def test_set_paragraph_writes_styled_run():
    prs   = Presentation()
    slide = prs.slides.add_slide( prs.slide_layouts[ 6 ] )
    tf    = slide.shapes.add_textbox( 0, 0, 100, 100 ).text_frame
    PptxDeckRenderer._set_paragraph(
        tf.paragraphs[ 0 ], "Hello", Pt( 20 ), RGBColor.from_string( "112233" ),
        bold=True, align=PP_ALIGN.CENTER,
    )
    para = tf.paragraphs[ 0 ]
    assert para.runs[ 0 ].text == "Hello"
    assert para.runs[ 0 ].font.bold is True
    assert para.alignment == PP_ALIGN.CENTER


# ---------------------------------------------------------------------------
# build — integration across slide types and branches
# ---------------------------------------------------------------------------

def _full_presentation():
    return PresentationModel(
        title="Owning Our Layout", speaker="Tiffany", date="2026-08-16",
        slides=[
            SlideModel(   # lead: title, with subtitle + speaker/date
                number=1, arc_position="opening", type="title",
                title="Owning Our Layout", subtitle="Real text at last",
                presenter_notes=PresenterNotes( talking_points=[ "Welcome" ], timing_seconds=30 ),
            ),
            SlideModel(   # lead: section_divider, no subtitle
                number=2, arc_position="body", type="section_divider",
                title="Part Two",
                presenter_notes=PresenterNotes( timing_seconds=0 ),   # → no notes slide
            ),
            SlideModel(   # content with subtitle + bullets + a real chart visual
                number=3, arc_position="body", type="content",
                title="Why It Matters", subtitle="Three wins",
                visual_type="chart", visual_description="bar chart",
                content_bullets=[ "Selectable", "Accessible", "Editable" ],
                presenter_notes=PresenterNotes( transition="So...", talking_points=[ "Explain" ] ),
            ),
            SlideModel(   # content, no subtitle, no bullets, no visual (all-False branches)
                number=4, arc_position="closing", type="conclusion",
                title="Thanks",
            ),
        ],
    )


def test_build_produces_real_text_layer( tmp_path ):
    visuals = tmp_path / "visuals"
    visuals.mkdir()
    _png( visuals / "chart-000.png" )             # slide 3 is the 1st non-text_only → idx 0
    out = str( tmp_path / "deck.pptx" )

    result = PptxDeckRenderer.build( _full_presentation(), {}, str( visuals ), out )

    assert result == out and os.path.isfile( out )
    verdict = verify_presentation_deck( out )
    assert verdict, verdict.reason
    assert verdict.slide_count == 4
    assert verdict.text_run_count > 0
    assert _slide_text_runs( out ) > 0


def test_build_embeds_visual_picture( tmp_path ):
    visuals = tmp_path / "visuals"
    visuals.mkdir()
    _png( visuals / "chart-000.png" )
    out = str( tmp_path / "deck.pptx" )
    PptxDeckRenderer.build( _full_presentation(), {}, str( visuals ), out )

    prs = Presentation( out )
    # slide index 2 (0-based) is the content slide with the chart.
    pictures = [ s for s in prs.slides[ 2 ].shapes if s.shape_type == 13 ]  # 13 = PICTURE
    assert len( pictures ) == 1


def test_build_content_slide_without_visual_file( tmp_path ):
    # Non-text_only slide but NO file on disk → no picture, text still real.
    pres = PresentationModel(
        title="No Visual", slides=[
            SlideModel(
                number=1, arc_position="body", type="content", title="Chartless",
                visual_type="chart", visual_description="missing",
                content_bullets=[ "still text" ],
            ),
        ],
    )
    out = str( tmp_path / "novisual.pptx" )
    PptxDeckRenderer.build( pres, {}, str( tmp_path / "visuals" ), out )
    prs = Presentation( out )
    pictures = [ s for s in prs.slides[ 0 ].shapes if s.shape_type == 13 ]
    assert pictures == []
    assert _slide_text_runs( out ) > 0


def test_build_empty_presentation( tmp_path ):
    out = str( tmp_path / "empty.pptx" )
    PptxDeckRenderer.build( PresentationModel( title="Empty" ), {}, str( tmp_path ), out )
    assert os.path.isfile( out )
    assert len( Presentation( out ).slides._sldIdLst ) == 0


def test_build_creates_missing_output_dir( tmp_path ):
    out = str( tmp_path / "nested" / "deep" / "deck.pptx" )
    PptxDeckRenderer.build(
        PresentationModel(
            title="T", slides=[ SlideModel( number=1, arc_position="opening", type="title", title="T" ) ],
        ),
        {}, str( tmp_path / "visuals" ), out,
    )
    assert os.path.isfile( out )


def test_build_notes_populated_and_omitted( tmp_path ):
    out = str( tmp_path / "notes.pptx" )
    PptxDeckRenderer.build( _full_presentation(), {}, str( tmp_path / "visuals" ), out )
    prs = Presentation( out )
    # Slide 1 has talking points → a notes slide with text.
    assert prs.slides[ 0 ].has_notes_slide
    assert "Welcome" in prs.slides[ 0 ].notes_slide.notes_text_frame.text
    # Slide 2 (section_divider, timing 0, no other note fields) → no notes slide.
    assert not prs.slides[ 1 ].has_notes_slide
